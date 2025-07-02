#!/usr/bin/env python3
import pandas as pd, matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

DF = pd.read_csv("simulation_results.csv").head(3)

# Radar Chart – simple revenue vs conv vs rsvp
labels = ["Revenue", "Conversion", "RSVP", "Retention"]
fig, ax = plt.subplots(subplot_kw={"projection": "polar"})
for _, row in DF.iterrows():
    data = [row["avg_gift_hkd"]*row["conv_rate"], row["conv_rate"], row["rsvp_pct"], row["retention_pct"], row["avg_gift_hkd"]*row["conv_rate"]]
    angles = [n/float(len(labels))*2*3.14159 for n in range(len(labels)+1)]
    ax.plot(angles, data, linewidth=1, linestyle='solid', label=f"Ask {row['ask']} | {row['format']}")
ax.set_xticks([n/float(len(labels))*2*3.14159 for n in range(len(labels))])
ax.set_xticklabels(labels)
plt.savefig("radar.png", bbox_inches='tight')

prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.add_picture("radar.png", Inches(1), Inches(1), width=Inches(8))
prs.save("kpi_report.pptx")
print("✅  Report → kpi_report.pptx")